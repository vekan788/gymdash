from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Helper to create slides
def add_slide(prs, title_text, content_items, img_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Content (No bullets)
    tf = slide.placeholders[1].text_frame
    tf.clear() # Clear default bullet
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
        # Remove bullet
        p.font.name = 'Arial'
    
    # Add Image (if exists and valid)
    if img_path:
        try:
            # Place image on the right
            left = Inches(5.5)
            top = Inches(2)
            height = Inches(4)
            slide.shapes.add_picture(img_path, left, top, height=height)
            
            # Adjust text width to not overlap
            slide.placeholders[1].width = Inches(4.5)
        except Exception as e:
            print(f"Skipping image {img_path}: {e}")

# --- SLIDE 1: Title ---
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Ein Tag im Leben eines Arztes"
subtitle.text = "Präsentiert von: Debre Nóra und Fekete-Kászoni Boróka Lea\nKlasse: XII.H"

# --- SLIDE 2: Aufgaben ---
add_slide(prs, "Aufgaben", [
    "untersuchen", 
    "die Diagnose", 
    "das Rezept", 
    "die Operation"
], "img_aufgaben.jpg")

# --- SLIDE 3: Ausbildung ---
add_slide(prs, "Ausbildung", [
    "das Medizinstudium", 
    "die Universität", 
    "sechs Jahre"
], "img_ausbildung.jpg") # Note: previous download failed (size 29 bytes), might skip or show placeholder if corrupted.

# --- SLIDE 4: Menschen & Materialien ---
add_slide(prs, "Menschen & Materialien", [
    "kranke Menschen", 
    "das Stethoskop", 
    "medizinische Geräte"
], "img_menschen.jpg")

# --- SLIDE 5: Eigenschaften ---
add_slide(prs, "Eigenschaften", [
    "geduldig", 
    "konzentriert", 
    "das Wissen"
], "img_eigenschaften.jpg")

# --- SLIDE 6: Arbeitsort ---
add_slide(prs, "Arbeitsort", [
    "das Krankenhaus", 
    "die Praxis"
], "img_arbeitsort.jpg")

# --- SLIDE 7: Arbeitszeiten ---
add_slide(prs, "Arbeitszeiten", [
    "lange Arbeitszeiten", 
    "der Schichtdienst", 
    "das Wochenende"
], "img_arbeitszeiten.jpg")

# --- SLIDE 8: Anstrengung ---
add_slide(prs, "Anstrengung", [
    "geistig anstrengend", 
    "die Verantwortung"
], "img_anstrengung.jpg")

# --- SLIDE 9: Vor- und Nachteile ---
add_slide(prs, "Vor- und Nachteile", [
    "helfen", 
    "das Gehalt", 
    "der Stress"
], "img_vorteile.jpg") # Note: download might have failed (29 bytes)

# --- SLIDE 10: Verdienst ---
add_slide(prs, "Verdienst", [
    "verdienen", 
    "brutto", 
    "steigen"
], "img_verdienst.jpg")

# --- SLIDE 11: Möglichkeiten ---
# Custom layout for this one as requested (beautifully put)
slide_layout = prs.slide_layouts[1] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Möglichkeiten"

# Add text box for centered beautiful text
left = Inches(1)
top = Inches(3)
width = Inches(8)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.clear()
p = tf.add_paragraph()
p.text = "der Oberarzt"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204) # Blueish

p2 = tf.add_paragraph()
p2.text = "die Forschung"
p2.alignment = PP_ALIGN.CENTER
p2.font.size = Pt(40)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0, 102, 204)

# --- SLIDE 12: Ende (Clean) ---
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Centered Thank You Text
left = Inches(1)
top = Inches(3)
width = Inches(8)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Vielen Dank für Ihre Aufmerksamkeit!"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(44)

# Save
prs.save("Arzt_Praesentation_v2.pptx")
print("Presentation generated successfully.")
