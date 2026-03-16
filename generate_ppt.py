from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Helper to create slides
def add_slide(prs, title_text, content_items):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Content (Bullets)
    tf = slide.placeholders[1].text_frame
    tf.text = content_items[0] if content_items else ""
    
    for item in content_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

# --- SLIDE 1: Title ---
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Ein Tag im Leben eines Arztes"
subtitle.text = "Präsentiert von: Debre Nóra und Fekete-Kászoni Boróka Lea\nKlasse: XII.H"

# --- SLIDE 2: Aufgaben ---
add_slide(prs, "Aufgaben", [
    "untersuchen", 
    "die Diagnose", 
    "das Rezept", 
    "die Operation"
])

# --- SLIDE 3: Ausbildung ---
add_slide(prs, "Ausbildung", [
    "das Medizinstudium", 
    "die Universität", 
    "sechs Jahre"
])

# --- SLIDE 4: Menschen & Materialien ---
add_slide(prs, "Menschen & Materialien", [
    "kranke Menschen", 
    "das Stethoskop", 
    "medizinische Geräte"
])

# --- SLIDE 5: Eigenschaften ---
add_slide(prs, "Eigenschaften", [
    "geduldig", 
    "konzentriert", 
    "das Wissen"
])

# --- SLIDE 6: Arbeitsort ---
add_slide(prs, "Arbeitsort", [
    "das Krankenhaus", 
    "die Praxis"
])

# --- SLIDE 7: Arbeitszeiten ---
add_slide(prs, "Arbeitszeiten", [
    "lange Arbeitszeiten", 
    "der Schichtdienst", 
    "das Wochenende"
])

# --- SLIDE 8: Anstrengung ---
add_slide(prs, "Anstrengung", [
    "geistig anstrengend", 
    "die Verantwortung"
])

# --- SLIDE 9: Vor- und Nachteile ---
add_slide(prs, "Vor- und Nachteile", [
    "helfen", 
    "das Gehalt", 
    "der Stress"
])

# --- SLIDE 10: Verdienst ---
add_slide(prs, "Verdienst", [
    "verdienen", 
    "brutto", 
    "steigen"
])

# --- SLIDE 11: Möglichkeiten ---
# Custom layout for this one as requested (beautifully put)
slide_layout = prs.slide_layouts[1] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Möglichkeiten"

# Add text box for centered beautiful text
left = Inches(1)
top = Inches(3)
width = Inches(8)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "der Oberarzt"
p = tf.add_paragraph()
p.text = "die Forschung"
p.alignment = PP_ALIGN.CENTER
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[1].font.size = Pt(40)
tf.paragraphs[0].font.bold = True
tf.paragraphs[1].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204) # Blueish
tf.paragraphs[1].font.color.rgb = RGBColor(0, 102, 204)

# --- SLIDE 12: Ende ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Ende"
tf = slide.placeholders[1].text_frame
tf.text = "Vielen Dank für Ihre Aufmerksamkeit!"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(44)

# Save
prs.save("Arzt_Praesentation.pptx")
print("Presentation generated successfully.")
